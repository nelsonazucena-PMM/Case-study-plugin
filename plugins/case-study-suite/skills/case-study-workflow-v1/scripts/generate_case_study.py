#!/usr/bin/env python3
"""
Single-slide TELUS Digital case study generator. Renders the layout
defined in references/Case Study Template.md §4:

  1. Sidebar: industry icon to the left of the industry name. The icon is
     loaded from the skill's icons/ folder when the industry matches a
     known mapping (Automotive, Banking and FinTech, Games, Healthcare,
     Media, Retail, Tech, Telecomms, Travel and Hospitality). The script
     reads the source PNG's native dimensions and scales width from the
     locked height so the icon's original aspect ratio is preserved. If
     no icon matches, the built-in line-art package mark is used as a
     fallback.
  2. Metadata grid: labels render bold BLACK, values render in medium
     SLATE gray.
  3. KPI accent bar: a 1/3 FOREST green top segment and a 2/3 GALENA
     gray bottom segment.
  4. Footer: embeds the TELUS Digital wordmark image
     (references/TELUS_Digital_EN_Hor_RGB_Blk_2025.png) in the bottom-left
     corner of the slide, with "Confidential" placed to its right. Falls
     back to a text wordmark if the PNG is not present.
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn


# ---------- Brand palette (AI Standards §9) ----------
FOREST      = RGBColor(0x00, 0x80, 0x4A)
TELUS_GREEN = RGBColor(0x66, 0xCC, 0x00)
OBSIDIAN    = RGBColor(0x22, 0x22, 0x20)
BLACK       = RGBColor(0x00, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
PEARL       = RGBColor(0xFC, 0xFD, 0xFB)
MOONSTONE   = RGBColor(0xF3, 0xF3, 0xED)
SIDEBAR_BG  = RGBColor(0xF2, 0xF3, 0xED)
MARBLE      = RGBColor(0xDE, 0xE0, 0xD9)
GALENA      = RGBColor(0xB6, 0xB6, 0xB1)
SLATE       = RGBColor(0x59, 0x59, 0x56)

# ---------- Geometry ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SIDEBAR_W = Inches(4.27)
SIDEBAR_X = SLIDE_W - SIDEBAR_W
MAIN_W    = SIDEBAR_X
MAIN_X    = Inches(0)

FOOTER_H  = Inches(0.35)
FOOTER_Y  = SLIDE_H - FOOTER_H

MAIN_MARGIN_L = Inches(0.55)
MAIN_MARGIN_R = Inches(0.40)
MAIN_MARGIN_T = Inches(0.45)

FONT_NAME    = "HN for TELUS SA"
FONT_DISPLAY = "HN for TELUS SA Display"

# ---------- Brand assets ----------
# TELUS Digital horizontal wordmark (English, RGB, black). Lives in the
# skill's references/ folder so it can be resolved relative to this script.
SKILL_ROOT = Path(__file__).resolve().parent.parent
LOGO_PATH = SKILL_ROOT / "references" / "TELUS_Digital_EN_Hor_RGB_Blk_2025.png"

# Industry icons (line-art black PNGs, native aspect ratio preserved at render).
# Lookup table: each key is a substring matched case-insensitively against the
# `industry` field from the JSON payload. Order matters — first match wins, so
# more specific aliases sit above more generic ones.
ICONS_DIR = SKILL_ROOT / "icons"
INDUSTRY_ICONS = [
    ("automotive",   "Automotive.png"),
    ("fintech",      "Banking and FinTech.png"),
    ("financial",    "Banking and FinTech.png"),
    ("banking",      "Banking and FinTech.png"),
    ("game",         "Games.png"),
    ("gaming",       "Games.png"),
    ("healthcare",   "Healthcare.png"),
    ("health",       "Healthcare.png"),
    ("pharma",       "Healthcare.png"),
    ("media",        "Media.png"),
    ("entertainment","Media.png"),
    ("streaming",    "Media.png"),
    ("retail",       "Retail.png"),
    ("e-commerce",   "Retail.png"),
    ("ecommerce",    "Retail.png"),
    ("telecom",      "Telecomms.png"),
    ("telco",        "Telecomms.png"),
    ("travel",       "Travel and Hospitality.png"),
    ("hospitality",  "Travel and Hospitality.png"),
    ("airline",      "Travel and Hospitality.png"),
    ("hotel",        "Travel and Hospitality.png"),
    # "tech" is intentionally last — it would otherwise capture "fintech".
    ("tech",         "Tech.png"),
    ("technology",   "Tech.png"),
    ("software",     "Tech.png"),
    ("saas",         "Tech.png"),
]


def _resolve_industry_icon(industry):
    """Return the Path to the icon PNG that matches the industry string,
    or None if no entry in INDUSTRY_ICONS matches. The match is a simple
    case-insensitive substring scan against the industry field."""
    if not industry:
        return None
    needle = industry.lower()
    for keyword, fname in INDUSTRY_ICONS:
        if keyword in needle:
            path = ICONS_DIR / fname
            if path.exists():
                return path
    return None


def _set_run(run, text, *, size_pt, bold=False, color=OBSIDIAN, font=FONT_NAME):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_rect(slide, x, y, w, h, rgb, *, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb
    sppr = shp._element.spPr
    for shd in sppr.findall(qn("a:effectLst")):
        sppr.remove(shd)
    return shp


def _add_text(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, word_wrap=True,
              auto_fit=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.paragraphs[0].text = ""
    return tb, tf


def _new_para(tf, first=False):
    if first and (tf.paragraphs[0].text == "" and not tf.paragraphs[0].runs):
        return tf.paragraphs[0]
    return tf.add_paragraph()


# ---------- Industry icon ----------

def _draw_package_icon(slide, x, y, size, color=BLACK, weight_pt=1.75):
    """Simple line-art package/box icon: outline rectangle + horizontal seam
    across the middle. Reads as 'shipping/logistics' without requiring an
    external image asset."""
    # Box body (outline only)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    box.fill.background()
    box.line.color.rgb = color
    box.line.width = Pt(weight_pt)
    sppr = box._element.spPr
    for shd in sppr.findall(qn("a:effectLst")):
        sppr.remove(shd)
    # Horizontal seam (tape line) across the middle
    seam_y = y + size // 2
    seam = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      x, seam_y, x + size, seam_y)
    seam.line.color.rgb = color
    seam.line.width = Pt(weight_pt)
    # Short vertical notch at the top center (flap)
    notch_x = x + size // 2
    notch = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       notch_x, y, notch_x, seam_y)
    notch.line.color.rgb = color
    notch.line.width = Pt(weight_pt)
    return box


# ---------- Section renderers ----------

def render_sidebar(slide, data):
    _add_rect(slide, SIDEBAR_X, Inches(0), SIDEBAR_W, SLIDE_H, SIDEBAR_BG)

    pad_x = Inches(0.35)
    content_w = SIDEBAR_W - pad_x * 2
    cursor_y = Inches(0.55)

    # --- Industry row: icon + name side by side (no uppercase micro-label) ---
    # Icon height is locked to the sidebar slot; width scales from the source
    # PNG's native dimensions so the line-art icon never stretches.
    icon_h = Inches(0.52)
    icon_x = SIDEBAR_X + pad_x
    icon_y = cursor_y
    icon_path = _resolve_industry_icon(data.get("industry", ""))

    if icon_path is not None:
        # python-pptx preserves aspect ratio when only one dimension is given.
        # Passing height alone scales width proportionally from the source PNG.
        pic = slide.shapes.add_picture(
            str(icon_path), icon_x, icon_y, height=icon_h,
        )
        icon_w = pic.width
    else:
        # Fallback: the built-in line-art package mark (square footprint).
        _draw_package_icon(slide, icon_x, icon_y, icon_h)
        icon_w = icon_h

    name_x = icon_x + icon_w + Inches(0.18)
    name_w = content_w - icon_w - Inches(0.18)
    tb, tf = _add_text(slide, name_x, icon_y, name_w, icon_h,
                       anchor=MSO_ANCHOR.MIDDLE)
    p = _new_para(tf, first=True)
    p.alignment = PP_ALIGN.LEFT
    _set_run(p.add_run(), data.get("industry", ""),
             size_pt=18, bold=True, color=BLACK)
    cursor_y += icon_h + Inches(0.06)

    sublabel = data.get("industry_sublabel")
    if sublabel:
        tb, tf = _add_text(slide, SIDEBAR_X + pad_x, cursor_y,
                           content_w, Inches(0.28))
        p = _new_para(tf, first=True)
        p.alignment = PP_ALIGN.LEFT
        _set_run(p.add_run(), sublabel, size_pt=11, color=SLATE)
        cursor_y += Inches(0.28)

    # Divider
    cursor_y += Inches(0.10)
    _add_rect(slide, SIDEBAR_X + pad_x, cursor_y, content_w, Emu(9525), MARBLE)
    cursor_y += Inches(0.18)

    # --- Metadata grid: labels BLACK bold, values SLATE gray ---
    metadata_pairs = [
        ("CHANNELS",          data.get("channels", "")),
        ("LANGUAGES",         data.get("languages", "")),
        ("DELIVERY GEOS",     data.get("delivery_geos", "")),
        ("PARTNERSHIP SINCE", data.get("partnership_since", "")),
    ]
    scope = data.get("scope")
    if scope:
        metadata_pairs.append(("SCOPE", scope))

    col_w = (content_w - Inches(0.15)) / 2
    row_h = Inches(0.60)
    for i, (label, value) in enumerate(metadata_pairs):
        col = i % 2
        row = i // 2
        cx = SIDEBAR_X + pad_x + col * (col_w + Inches(0.15))
        cy = cursor_y + row * row_h
        tb, tf = _add_text(slide, cx, cy, col_w, row_h)
        p = _new_para(tf, first=True)
        _set_run(p.add_run(), label, size_pt=11, bold=True, color=BLACK)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        _set_run(p2.add_run(), value, size_pt=11, color=SLATE)

    rows = (len(metadata_pairs) + 1) // 2
    cursor_y += row_h * rows + Inches(0.10)

    # Divider
    _add_rect(slide, SIDEBAR_X + pad_x, cursor_y, content_w, Emu(6350), MARBLE)
    cursor_y += Inches(0.15)

    # Photo placeholder
    photo_y = cursor_y
    photo_h = SLIDE_H - photo_y - FOOTER_H - Inches(0.05)
    photo = _add_rect(slide, SIDEBAR_X, photo_y, SIDEBAR_W, photo_h, GALENA)
    tb, tf = _add_text(slide, SIDEBAR_X, photo_y, SIDEBAR_W, photo_h,
                       anchor=MSO_ANCHOR.MIDDLE)
    p = _new_para(tf, first=True)
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), "[ Insert industry-contextual photo ]",
             size_pt=10, color=WHITE, bold=True)


def render_title(slide, data):
    title = data.get("title", "")
    accent = data.get("title_accent", "").strip()
    tb, tf = _add_text(
        slide, MAIN_MARGIN_L, Inches(0.35),
        MAIN_W - MAIN_MARGIN_L - MAIN_MARGIN_R, Inches(0.95),
        auto_fit=True,
    )
    p = _new_para(tf, first=True)
    p.alignment = PP_ALIGN.LEFT
    if accent and accent in title:
        before, _, after = title.partition(accent)
        if before:
            _set_run(p.add_run(), before, size_pt=30, bold=False, color=BLACK, font=FONT_DISPLAY)
        _set_run(p.add_run(), accent, size_pt=30, bold=True, color=FOREST, font=FONT_DISPLAY)
        if after:
            _set_run(p.add_run(), after, size_pt=30, bold=False, color=BLACK, font=FONT_DISPLAY)
    else:
        _set_run(p.add_run(), title, size_pt=30, bold=True, color=BLACK, font=FONT_DISPLAY)


def render_challenge_outcome(slide, data):
    top_y = Inches(1.45)
    col_w = (MAIN_W - MAIN_MARGIN_L - MAIN_MARGIN_R - Inches(0.35)) / 2
    col_h = Inches(2.50)
    for i, (header, body, color) in enumerate([
        ("Challenge", data.get("challenge", ""), BLACK),
        ("Outcome",   data.get("outcome", ""),   BLACK),
    ]):
        x = MAIN_MARGIN_L + i * (col_w + Inches(0.35))
        tb, tf = _add_text(slide, x, top_y, col_w, col_h, auto_fit=True)
        p = _new_para(tf, first=True)
        _set_run(p.add_run(), header, size_pt=14, bold=True, color=color)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        _set_run(p2.add_run(), body, size_pt=10.5, color=OBSIDIAN)


def render_solution(slide, data):
    top_y = Inches(4.05)
    w = MAIN_W - MAIN_MARGIN_L - MAIN_MARGIN_R
    h = Inches(1.90)
    tb, tf = _add_text(slide, MAIN_MARGIN_L, top_y, w, h, auto_fit=True)
    p = _new_para(tf, first=True)
    _set_run(p.add_run(), "Solution", size_pt=14, bold=True, color=FOREST)
    for paragraph in (data.get("solution", "") or "").split("\n\n"):
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        _set_run(p2.add_run(), paragraph, size_pt=10.5, color=OBSIDIAN)


def render_kpi_bar(slide, data):
    kpis = data.get("kpis", [])
    n = len(kpis)
    if n < 3 or n > 4:
        raise ValueError(f"KPI count must be 3 or 4, got {n}")

    top_y = Inches(6.05)
    bar_h = Inches(1.00)
    total_w = MAIN_W - MAIN_MARGIN_L - MAIN_MARGIN_R
    gap = Inches(0.20)
    cell_w = (total_w - gap * (n - 1)) / n
    accent_w = Emu(25400)

    for i, kpi in enumerate(kpis):
        cx = MAIN_MARGIN_L + i * (cell_w + gap)

        # --- Split accent bar: top 1/3 FOREST green, bottom 2/3 GALENA gray ---
        green_h = bar_h // 3
        gray_h = bar_h - green_h
        _add_rect(slide, cx, top_y, accent_w, green_h, FOREST)
        _add_rect(slide, cx, top_y + green_h, accent_w, gray_h, GALENA)

        tb, tf = _add_text(slide, cx + Inches(0.10), top_y,
                           cell_w - Inches(0.10), bar_h,
                           anchor=MSO_ANCHOR.TOP)
        p = _new_para(tf, first=True)
        _set_run(p.add_run(), kpi.get("number", ""),
                 size_pt=26, bold=True, color=FOREST, font=FONT_DISPLAY)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        _set_run(p2.add_run(), kpi.get("label", ""),
                 size_pt=11, bold=True, color=OBSIDIAN)
        context = (kpi.get("context") or "").strip()
        if context:
            p3 = tf.add_paragraph()
            p3.space_before = Pt(1)
            _set_run(p3.add_run(), context, size_pt=9, color=SLATE)


def render_footer(slide):
    sep_y = FOOTER_Y - Inches(0.04)
    _add_rect(slide, MAIN_MARGIN_L, sep_y,
              MAIN_W - MAIN_MARGIN_L - MAIN_MARGIN_R,
              Emu(6350), MARBLE)

    # --- TELUS Digital logo, bottom-left corner of the slide ---
    # The PNG aspect ratio is roughly 4:1 (wordmark + symbol). Lock the
    # height to fit inside the footer band and let pptx scale the width
    # proportionally so the mark never distorts.
    logo_h = Inches(0.26)
    logo_y = FOOTER_Y + (FOOTER_H - logo_h) / 2
    logo_x = MAIN_MARGIN_L
    logo_right_edge = logo_x  # fallback if image is missing

    if LOGO_PATH.exists():
        pic = slide.shapes.add_picture(
            str(LOGO_PATH), logo_x, logo_y, height=logo_h,
        )
        logo_right_edge = pic.left + pic.width
    else:
        # Graceful fallback: if the asset is missing, render the wordmark
        # as text so the footer is never empty.
        tb, tf = _add_text(slide, logo_x, FOOTER_Y, Inches(1.4), FOOTER_H,
                           anchor=MSO_ANCHOR.MIDDLE)
        p = _new_para(tf, first=True)
        _set_run(p.add_run(), "TELUS Digital",
                 size_pt=9, bold=True, color=BLACK)
        logo_right_edge = logo_x + Inches(1.4)

    # --- "Confidential" sits to the right of the logo on the same baseline ---
    conf_x = logo_right_edge + Inches(0.18)
    conf_w = (MAIN_X + MAIN_W - MAIN_MARGIN_R) - conf_x
    if conf_w > Inches(0.5):
        tb, tf = _add_text(slide, conf_x, FOOTER_Y, conf_w, FOOTER_H,
                           anchor=MSO_ANCHOR.MIDDLE)
        p = _new_para(tf, first=True)
        p.alignment = PP_ALIGN.LEFT
        _set_run(p.add_run(), "Confidential", size_pt=8, color=SLATE)


def set_speaker_notes(slide, data):
    notes = data.get("speaker_notes", {})
    lines = [
        f"Client: {notes.get('client_name', '[REDACTED]')} (Please remove before sharing externally)",
        f"Service Line(s): {notes.get('service_lines', '')}",
        f"Industry: {notes.get('industry_full', '')}",
        f"Logo Permission: {notes.get('logo_permission', 'No')}",
    ]
    if notes.get("ops_leads"):
        lines.append(f"Ops Leads: {notes['ops_leads']}")
    nf = slide.notes_slide.notes_text_frame
    nf.clear()
    nf.text = lines[0]
    for line in lines[1:]:
        p = nf.add_paragraph()
        p.text = line


def build_presentation(data, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, Inches(0), Inches(0), MAIN_W, SLIDE_H, PEARL)
    render_sidebar(slide, data)
    render_title(slide, data)
    render_challenge_outcome(slide, data)
    render_solution(slide, data)
    render_kpi_bar(slide, data)
    render_footer(slide)
    set_speaker_notes(slide, data)
    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) != 3:
        print("Usage: python3 generate_case_study.py <input.json> <output.pptx>",
              file=sys.stderr)
        sys.exit(2)
    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])
    data = json.loads(input_path.read_text(encoding="utf-8"))
    required_top = ["title", "industry", "channels", "languages",
                    "delivery_geos", "partnership_since",
                    "challenge", "solution", "outcome", "kpis", "speaker_notes"]
    missing = [k for k in required_top if k not in data]
    if missing:
        raise ValueError(f"Missing required fields: {missing}")
    build_presentation(data, output_path)
    print(f"Wrote {output_path}")


if __name__ == "__main__":
    main()
